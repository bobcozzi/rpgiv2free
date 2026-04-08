
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

def add_code_slide(title, before, after):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = "Before:\n" + before + "\n\nAfter:\n" + after

# Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Modern RPG IV (7.3 → 7.6)"
slide.placeholders[1].text = "High-impact enhancements developers should actually use"

# DATA-INTO
add_slide("DATA-INTO (7.3)",
          ["Parse JSON/XML directly",
           "Replaces manual parsing",
           "Cleaner, faster code"])

add_code_slide("DATA-INTO Example",
"""// Old parsing
read jsonLine;
dow not %eof();
  // manual parsing logic
enddo;""",
"""data-into myDS %data(jsonDoc)
  %parser('YAJL/YAJLINTO');""")

# FOR-EACH
add_slide("FOR-EACH (7.4 / 7.3 PTF)",
          ["Simplifies loops",
           "Replaces READ/DO pattern",
           "More readable"])

add_code_slide("FOR-EACH Example",
"""dou %eof(file);
  read file;
enddo;""",
"""for-each rec in file;
endfor;""")

# %SPLIT
add_slide("%SPLIT (7.4 TR / 7.3 PTF)",
          ["Break strings into arrays",
           "Eliminates custom parsing"])

add_code_slide("%SPLIT Example",
"""// manual parsing
pos = %scan(',' : str);""",
"""arr = %split(str : ',');""")

# %LOWER / %UPPER
add_slide("%LOWER / %UPPER",
          ["Case conversion",
           "Cleaner than custom logic"])

add_code_slide("%LOWER Example",
"""// old
for i = 1 to %len(str);""",
"""str = %lower(str);""")

# %LIST / %RANGE
add_slide("%LIST / %RANGE",
          ["Inline lists and comparisons",
           "Cleaner conditionals"])

add_code_slide("%LIST Example",
"""if code = 'A' or code = 'B' or code = 'C';""",
"""if code in %list('A':'B':'C');""")

# ENUM
add_slide("ENUM (7.6)",
          ["Stronger typing",
           "Prevents invalid values"])

# SND-MSG
add_slide("SND-MSG",
          ["Simplified messaging",
           "Cleaner than APIs"])

# %HIVAL / %LOVAL
add_slide("%HIVAL / %LOVAL (7.6)",
          ["Boundary values",
           "Cleaner than hardcoding"])

# Summary
add_slide("Use These First",
          ["DATA-INTO",
           "FOR-EACH",
           "%SPLIT",
           "%LOWER/%UPPER",
           "%LIST/%RANGE"])

path = "/Users/cozzi/Downloads/RPGIV_SINCE.pptx"
prs.save(path)
