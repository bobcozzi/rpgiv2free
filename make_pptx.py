#!/usr/bin/env python3
"""
make_pptx.py
Converts VSCODE_RPG_PRESENTATION.md -> VSCODE_RPG_PRESENTATION.pptx
  • One slide per ## Slide N section
  • Code blocks: Letter Gothic 10pt, dark background
  • Tables: styled header row + alternating rows
  • Bold / italic / inline-code inline formatting preserved
"""

import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

MD_FILE      = sys.argv[1] if len(sys.argv) > 1 else "VSCODE_RPG_PRESENTATION.md"
OUT_FILE     = sys.argv[2] if len(sys.argv) > 2 else "VSCODE_RPG_PRESENTATION_v2.pptx"
TEMPLATE     = "wisp_template.pptx"

# ── Dimensions ────────────────────────────────────────────────────────────────
W        = Inches(13.333)
H        = Inches(7.5)
MARGIN   = Inches(0.45)
TITLE_H  = Inches(1.05)
BODY_TOP = Inches(1.65)
BODY_L   = Inches(0.5)
BODY_W   = Inches(12.3)

# ── Colour palette ────────────────────────────────────────────────────────────
C_TITLE_BG = RGBColor(0x1F, 0x49, 0x7D)   # Navy blue
C_TITLE_FG = RGBColor(0xFF, 0xFF, 0xFF)
C_SLIDE_BG = RGBColor(0xF5, 0xF7, 0xFA)   # Near-white
C_BODY     = RGBColor(0x1A, 0x1A, 0x2E)
C_BOLD     = RGBColor(0x1F, 0x49, 0x7D)
C_NOTE     = RGBColor(0x66, 0x66, 0x66)
C_CODE_BG  = RGBColor(0x1E, 0x1E, 0x1E)   # VSCode dark
C_CODE_FG  = RGBColor(0xD4, 0xD4, 0xD4)
C_INLINE   = RGBColor(0xB5, 0x20, 0x4F)   # Inline `code`
C_TBL_HDR  = RGBColor(0x1F, 0x49, 0x7D)
C_TBL_ALT  = RGBColor(0xEA, 0xF0, 0xFA)

BULLET_CHARS = {0: '•', 1: '◦', 2: '▪'}


# ── Utility ───────────────────────────────────────────────────────────────────

def plain(text):
    """Strip **/ */ ` markers from text."""
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*',     r'\1', text)
    text = re.sub(r'`(.+?)`',       r'\1', text)
    return text


def set_box_margins(tf, l=Inches(0.12), r=Inches(0.12),
                    t=Inches(0.10), b=Inches(0.10)):
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('lIns', str(int(l)))
        bp.set('rIns', str(int(r)))
        bp.set('tIns', str(int(t)))
        bp.set('bIns', str(int(b)))


# ── Markdown parser ───────────────────────────────────────────────────────────

def parse_slides(md):
    """Return list of {title, blocks} dicts, one per slide section."""
    # Strip YAML / Marp frontmatter (--- ... ---)
    md = re.sub(r'^---\n.*?\n---\n', '', md, flags=re.DOTALL)

    sections = re.split(r'\n---\n', md)
    slides = []
    for sec in sections:
        sec = sec.strip()
        if not sec:
            continue
        m = re.match(r'^(#{1,3}) (.+)$', sec, re.MULTILINE)
        if not m:
            continue
        title = plain(m.group(2).strip())
        # Remove "Slide N — " prefix so it reads cleanly
        title = re.sub(r'^Slide\s+\d+\s*[—\-]\s*', '', title)
        rest   = sec[m.end():].strip()
        blocks = parse_blocks(rest)
        # Skip preamble sections that have no content blocks
        if not blocks:
            continue
        slides.append({'title': title, 'blocks': blocks})
    return slides


def parse_blocks(text):
    """Parse content text into a list of typed block dicts."""
    blocks = []
    lines  = text.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]

        # ── Code fence ──────────────────────────────────────────────────────
        if re.match(r'^\s*```', line):
            code = []
            i += 1
            while i < len(lines) and not re.match(r'^\s*```', lines[i]):
                code.append(lines[i])
                i += 1
            blocks.append({'type': 'code', 'text': '\n'.join(code)})
            i += 1
            continue

        # ── Table ────────────────────────────────────────────────────────────
        if line.strip().startswith('|'):
            rows = []
            while i < len(lines) and lines[i].strip().startswith('|'):
                row = lines[i].strip()
                # Skip separator rows like |---|---|
                if not re.fullmatch(r'[\|\s:\-]+', row):
                    rows.append([plain(c.strip()) for c in row.strip('|').split('|')])
                i += 1
            if rows:
                blocks.append({'type': 'table', 'rows': rows})
            continue

        # ── Blockquote ───────────────────────────────────────────────────────
        if line.strip().startswith('>'):
            txt = re.sub(r'^>\s*', '', line.strip())
            blocks.append({'type': 'quote', 'text': txt})
            i += 1
            continue

        # ── Bullet (- item, with optional leading spaces for nesting) ────────
        bm = re.match(r'^(\s*)-\s+(.+)$', line)
        if bm:
            level = min(len(bm.group(1)) // 2, 2)
            blocks.append({'type': 'bullet', 'text': bm.group(2), 'level': level})
            i += 1
            continue

        # ── Italic-only line  *text* ──────────────────────────────────────────
        if re.match(r'^\*[^*].+[^*]\*$', line.strip()):
            blocks.append({'type': 'italic', 'text': line.strip()[1:-1]})
            i += 1
            continue

        # ── Sub-heading inside content (### ...) → treat as bold para ─────────
        hm = re.match(r'^#{2,4}\s+(.+)$', line.strip())
        if hm:
            blocks.append({'type': 'subhead', 'text': plain(hm.group(1))})
            i += 1
            continue

        # ── Empty / skip ──────────────────────────────────────────────────────
        if not line.strip():
            i += 1
            continue

        # ── Plain paragraph ────────────────────────────────────────────────────
        blocks.append({'type': 'para', 'text': line.strip()})
        i += 1
    return blocks


# ── Rich-text run helpers ─────────────────────────────────────────────────────

def add_run(para, text, size=Pt(15), bold=False, italic=False,
            color=None, font='Calibri'):
    run = para.add_run()
    run.text        = text
    run.font.name   = font
    run.font.size   = size
    run.font.bold   = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return run


def add_inline(para, text, size=Pt(15), default_color=None):
    """Parse **bold**, *italic*, `code` inline markers and emit runs."""
    if default_color is None:
        default_color = C_BODY
    tokens = re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`)', text)
    for tok in tokens:
        if not tok:
            continue
        if tok.startswith('**') and tok.endswith('**'):
            add_run(para, tok[2:-2], size=size, bold=True, color=C_BOLD)
        elif tok.startswith('*') and tok.endswith('*'):
            add_run(para, tok[1:-1], size=size, italic=True, color=C_NOTE)
        elif tok.startswith('`') and tok.endswith('`'):
            add_run(para, tok[1:-1], size=Pt(13), font='Courier New', color=C_INLINE)
        else:
            add_run(para, tok, size=size, color=default_color)


# ── Shape builders ────────────────────────────────────────────────────────────

def get_ph(slide, idx):
    """Return placeholder by idx (0=title, 1=body), or None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def add_title_bar(slide, title_text):
    """Write text into the title placeholder; Wisp theme handles the styling."""
    ph = get_ph(slide, 0)
    if ph is None:
        return
    tf = ph.text_frame
    tf.word_wrap = True
    tf.clear()
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text


def estimate_shape_h(blocks):
    """Rough height needed by code/table blocks only (used to size body placeholder)."""
    total = Inches(0)
    for b in blocks:
        if b['type'] == 'code':
            n = max(len(b['text'].splitlines()), 1)
            total += Inches(0.265 * n + 0.28) + Inches(0.14)
        elif b['type'] == 'table':
            total += Inches(0.44 * len(b['rows'])) + Inches(0.14)
        else:
            txt = plain(b.get('text', ''))
            total += Inches(0.36 * max(1, len(txt) // 72 + 1) + 0.06)
    return total


def fill_body_ph(slide, text_blocks, top, height):
    """Write text/bullet blocks into the real body placeholder (idx=1). Returns next top."""
    if not text_blocks:
        return top
    ph = get_ph(slide, 1)
    if ph is None:
        # No body placeholder on this layout — fall back to a plain textbox
        box = slide.shapes.add_textbox(BODY_L, top, BODY_W, height)
        ph  = box
    ph.left   = BODY_L
    ph.top    = top
    ph.width  = BODY_W
    ph.height = height
    tf = ph.text_frame
    tf.word_wrap = True
    set_box_margins(tf)
    tf.clear()
    first = True
    for blk in text_blocks:
        p     = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        btype = blk['type']
        if btype == 'bullet':
            p.level = blk.get('level', 0)
            add_inline(p, blk['text'], size=Pt(15))
            p.space_before = Pt(2)
            p.space_after  = Pt(2)
        elif btype == 'quote':
            p.level = 0
            add_run(p, blk['text'], size=Pt(14), italic=True, color=C_NOTE)
            p.space_before = Pt(6)
            p.space_after  = Pt(4)
        elif btype == 'italic':
            p.level = 0
            add_run(p, blk['text'], size=Pt(15), italic=True, color=C_NOTE)
            p.space_before = Pt(4)
        elif btype == 'subhead':
            p.level = 0
            add_run(p, blk['text'], size=Pt(16), bold=True, color=C_BOLD)
            p.space_before = Pt(8)
            p.space_after  = Pt(2)
        else:
            p.level = 0
            add_inline(p, blk['text'], size=Pt(15))
            p.space_before = Pt(4)
            p.space_after  = Pt(2)
    return top + height + Inches(0.1)


def build_code_box(slide, code_text, top):
    """Emit a dark code box with Letter Gothic 10pt. Returns new top."""
    lines  = code_text.splitlines()
    nlines = max(len(lines), 1)
    h      = max(Inches(0.265 * nlines + 0.28), Inches(0.5))

    box = slide.shapes.add_textbox(BODY_L, top, BODY_W, h)
    box.fill.solid()
    box.fill.fore_color.rgb = C_CODE_BG
    tf = box.text_frame
    tf.word_wrap = False
    set_box_margins(tf, l=Inches(0.16), r=Inches(0.12),
                    t=Inches(0.14), b=Inches(0.12))

    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run               = p.add_run()
        run.text          = line if line.strip() else ' '
        run.font.name     = 'Letter Gothic MT'
        run.font.size     = Pt(10)
        run.font.color.rgb = C_CODE_FG
        run.font.bold     = False

    return top + h + Inches(0.14)


def build_table(slide, rows, top):
    """Emit a styled PPTX table. Returns new top."""
    if not rows:
        return top
    ncols  = max(len(r) for r in rows)
    nrows  = len(rows)
    row_h  = Inches(0.44)
    tbl_h  = row_h * nrows

    tbl = slide.shapes.add_table(
        nrows, ncols, BODY_L, top, BODY_W, tbl_h
    ).table

    for ri, row in enumerate(rows):
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            val  = row[ci] if ci < len(row) else ''
            tf   = cell.text_frame
            tf.word_wrap = True
            p    = tf.paragraphs[0]
            run  = p.add_run()
            run.text          = val
            run.font.name     = 'Calibri'
            run.font.size     = Pt(12)
            if ri == 0:
                run.font.bold      = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_TBL_HDR
            else:
                run.font.color.rgb = C_BODY
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C_TBL_ALT

    return top + tbl_h + Inches(0.14)


def add_content(slide, blocks):
    """Lay out all blocks on the slide.

    Text/bullet content before the first code or table goes into the real body
    placeholder (idx=1) so PowerPoint treats it as native content.
    Code blocks and tables are placed as shapes below.
    """
    SLIDE_BOTTOM = H - Inches(0.25)
    BODY_AREA_H  = SLIDE_BOTTOM - BODY_TOP

    # Split: text that precedes the first code/table vs everything after
    first_special = next(
        (i for i, b in enumerate(blocks) if b['type'] in ('code', 'table')),
        len(blocks)
    )
    pre_blocks  = blocks[:first_special]   # -> body placeholder
    post_blocks = blocks[first_special:]   # -> shapes below placeholder

    # Work out how much vertical space the post-blocks need
    post_h = estimate_shape_h(post_blocks)

    if pre_blocks:
        body_h  = max(BODY_AREA_H - post_h - Inches(0.1), Inches(0.6))
        cur_top = fill_body_ph(slide, pre_blocks, BODY_TOP, body_h)
    else:
        # No text — hide the body placeholder so it doesn't show on screen
        ph = get_ph(slide, 1)
        if ph:
            ph.top    = Inches(8)
            ph.height = Inches(0.01)
        cur_top = BODY_TOP

    for blk in post_blocks:
        if blk['type'] == 'code':
            cur_top = build_code_box(slide, blk['text'], cur_top)
        elif blk['type'] == 'table':
            cur_top = build_table(slide, blk['rows'], cur_top)
        else:
            # Text block interspersed with code — small label textbox
            txt = plain(blk.get('text', ''))
            h   = Inches(0.36)
            box = slide.shapes.add_textbox(BODY_L, cur_top, BODY_W, h)
            tf  = box.text_frame
            tf.word_wrap = True
            p   = tf.paragraphs[0]
            add_inline(p, blk.get('text', txt), size=Pt(14))
            cur_top += h + Inches(0.06)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    md_path = Path(MD_FILE)
    if not md_path.exists():
        print(f"ERROR: {MD_FILE} not found")
        return

    md     = md_path.read_text(encoding='utf-8')
    slides = parse_slides(md)
    print(f"Parsed {len(slides)} slides from {MD_FILE}")

    import os
    template = TEMPLATE if os.path.exists(TEMPLATE) else None
    prs = Presentation(template)
    prs.slide_width  = W
    prs.slide_height = H

    # Use 'Title and Content' -- provides title (idx=0) and body (idx=1) placeholders
    tc_layout = next(
        (l for l in prs.slide_layouts if l.name == 'Title and Content'),
        prs.slide_layouts[1]
    )

    for i, sd in enumerate(slides, 1):
        slide = prs.slides.add_slide(tc_layout)
        add_title_bar(slide, sd['title'])
        add_content(slide, sd['blocks'])
        print(f"  Slide {i:2d}: {sd['title'][:60]}")

    prs.save(OUT_FILE)
    print(f"\nSaved → {OUT_FILE}")


if __name__ == '__main__':
    main()
